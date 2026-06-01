from app.engines.knowledge_engine.parser.base import BaseDocumentParser


class PPTParser(BaseDocumentParser):
    SUPPORTED_EXTENSIONS = [".pptx"]

    def parse(self, file_path: str, **kwargs) -> dict:
        from pptx import Presentation

        prs = Presentation(file_path)
        slides_text = []
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            slides_text.append("\n".join(slide_text))
        return {
            "content": "\n---\n".join(slides_text),
            "slides": slides_text,
            "metadata": {"format": "pptx", "slide_count": len(slides_text)},
        }
